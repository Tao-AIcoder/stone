"""
tools/office_tool.py - Office file operations tool for STONE (默行者)

Supports .docx (Word), .xlsx (Excel), .pptx (PowerPoint).
Phase 1b scope: content + basic styles (heading/bold/italic/table/cell format).
Images, formulas, animations: out of scope.

All files are sandboxed within WORKSPACE_DIR.
"""

from __future__ import annotations

import logging
import re
from pathlib import Path
from typing import Any

from config import settings
from tools.base import ToolInterface, ToolResult

logger = logging.getLogger(__name__)

_VALID_ACTIONS = ("create", "read", "append", "delete")
_SUPPORTED_FORMATS = (".docx", ".xlsx", ".pptx")


def _resolve_path(rel_path: str) -> Path | None:
    """Resolve path within WORKSPACE_DIR. Returns None if path traversal detected."""
    workspace = Path(getattr(settings, "workspace_dir", "workspace")).resolve()
    target = (workspace / rel_path).resolve()
    try:
        target.relative_to(workspace)
        return target
    except ValueError:
        return None


class OfficeTool(ToolInterface):
    """
    Create, read, append to, and delete Office documents.

    Actions:
      create  : create a new document with initial content + styles
      read    : extract text content (tables as Markdown)
      append  : add content to an existing document
      delete  : remove a file (requires confirmation)
    """

    name = "office_tool"
    description = (
        "创建、读取、编辑 Office 文档（.docx Word、.xlsx Excel、.pptx PPT）。"
        "支持标题、正文、粗斜体、列表、表格等基础格式。"
    )
    requires_confirmation = False

    def needs_confirmation_for(self, params: dict) -> bool:
        return params.get("action") == "delete"

    async def execute(
        self,
        params: dict[str, Any],
        user_id: str = "default_user",
    ) -> ToolResult:
        action = params.get("action", "")
        if action not in _VALID_ACTIONS:
            return ToolResult.fail(
                f"不支持的 action: {action!r}。"
                f"合法 action: {', '.join(_VALID_ACTIONS)}"
            )
        path_str = params.get("path", "").strip()
        if not path_str:
            return ToolResult.fail("path 不能为空")

        path = _resolve_path(path_str)
        if not path:
            return ToolResult.fail("路径不合法（不允许访问工作目录外的文件）")

        ext = path.suffix.lower()
        if ext not in _SUPPORTED_FORMATS:
            return ToolResult.fail(
                f"不支持的文件格式: {ext!r}。"
                f"支持: {', '.join(_SUPPORTED_FORMATS)}"
            )

        handler = getattr(self, f"_action_{action}")
        return await handler(params, path, ext)

    # ── Create ────────────────────────────────────────────────────────────────

    async def _action_create(
        self, params: dict[str, Any], path: Path, ext: str
    ) -> ToolResult:
        if path.exists():
            return ToolResult.fail(f"文件已存在: {path.name}。如要覆盖请先 delete。")
        path.parent.mkdir(parents=True, exist_ok=True)
        content = params.get("content", "")
        title = params.get("title", path.stem)
        styles = params.get("styles", {})

        try:
            if ext == ".docx":
                self._create_docx(path, title, content, styles)
            elif ext == ".xlsx":
                self._create_xlsx(path, title, content, styles)
            elif ext == ".pptx":
                self._create_pptx(path, title, content, styles)
            return ToolResult.ok(
                f"已创建: {path.name}",
                {"path": str(path), "format": ext},
            )
        except ImportError as exc:
            return ToolResult.fail(f"缺少依赖库: {exc}。请执行 pip install python-docx openpyxl python-pptx")
        except Exception as exc:
            return ToolResult.fail(f"创建失败: {exc}")

    # ── Read ──────────────────────────────────────────────────────────────────

    async def _action_read(
        self, params: dict[str, Any], path: Path, ext: str
    ) -> ToolResult:
        if not path.exists():
            return ToolResult.fail(f"文件不存在: {path.name}")
        try:
            if ext == ".docx":
                content = self._read_docx(path)
            elif ext == ".xlsx":
                content = self._read_xlsx(path)
            elif ext == ".pptx":
                content = self._read_pptx(path)
            else:
                return ToolResult.fail(f"不支持的格式: {ext}")
            return ToolResult.ok(content, {"path": str(path)})
        except ImportError as exc:
            return ToolResult.fail(f"缺少依赖库: {exc}")
        except Exception as exc:
            return ToolResult.fail(f"读取失败: {exc}")

    # ── Append ────────────────────────────────────────────────────────────────

    async def _action_append(
        self, params: dict[str, Any], path: Path, ext: str
    ) -> ToolResult:
        if not path.exists():
            return ToolResult.fail(f"文件不存在: {path.name}，请先 create")
        content = params.get("content", "").strip()
        if not content:
            return ToolResult.fail("content 不能为空")
        styles = params.get("styles", {})
        try:
            if ext == ".docx":
                self._append_docx(path, content, styles)
            elif ext == ".xlsx":
                self._append_xlsx(path, content)
            elif ext == ".pptx":
                self._append_pptx(path, content)
            return ToolResult.ok(f"已追加内容到: {path.name}")
        except ImportError as exc:
            return ToolResult.fail(f"缺少依赖库: {exc}")
        except Exception as exc:
            return ToolResult.fail(f"追加失败: {exc}")

    # ── Delete ────────────────────────────────────────────────────────────────

    async def _action_delete(
        self, params: dict[str, Any], path: Path, ext: str
    ) -> ToolResult:
        if not path.exists():
            return ToolResult.fail(f"文件不存在: {path.name}")
        path.unlink()
        return ToolResult.ok(f"已删除: {path.name}")

    # ── Word (.docx) ──────────────────────────────────────────────────────────

    def _create_docx(self, path: Path, title: str, content: str, styles: dict) -> None:
        from docx import Document  # type: ignore[import]
        from docx.shared import Pt  # type: ignore[import]
        doc = Document()
        if title:
            doc.add_heading(title, level=1)
        if content:
            self._docx_add_content(doc, content, styles)
        doc.save(path)

    def _append_docx(self, path: Path, content: str, styles: dict) -> None:
        from docx import Document  # type: ignore[import]
        doc = Document(path)
        self._docx_add_content(doc, content, styles)
        doc.save(path)

    def _docx_add_content(self, doc: Any, content: str, styles: dict) -> None:
        """Parse Markdown-like content and add to docx."""
        from docx.shared import Pt  # type: ignore[import]
        bold = styles.get("bold", False)
        italic = styles.get("italic", False)
        font_size = styles.get("font_size")
        heading_level = styles.get("heading_level")

        for line in content.split("\n"):
            # Heading detection
            heading_match = re.match(r"^(#{1,6})\s+(.*)", line)
            if heading_match:
                level = len(heading_match.group(1))
                doc.add_heading(heading_match.group(2), level=min(level, 9))
                continue
            # Table row detection
            if "|" in line and line.strip().startswith("|"):
                # handled as part of table block below (simplified single-row)
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                table = doc.add_table(rows=1, cols=len(cells))
                for i, cell in enumerate(cells):
                    table.cell(0, i).text = cell
                continue
            # Separator line
            if re.match(r"^[-|=]{3,}$", line.strip()):
                continue
            # Normal paragraph
            if line.strip():
                para = doc.add_paragraph()
                run = para.add_run(line)
                if bold or styles.get("bold"):
                    run.bold = True
                if italic or styles.get("italic"):
                    run.italic = True
                if font_size:
                    run.font.size = Pt(int(font_size))
            elif heading_level:
                doc.add_heading(content, level=int(heading_level))

    def _read_docx(self, path: Path) -> str:
        from docx import Document  # type: ignore[import]
        doc = Document(path)
        parts: list[str] = []
        for para in doc.paragraphs:
            style = para.style.name if para.style else ""
            if style.startswith("Heading"):
                level = re.search(r"\d", style)
                hashes = "#" * int(level.group()) if level else "#"
                parts.append(f"{hashes} {para.text}")
            elif para.text.strip():
                parts.append(para.text)
        # Tables
        for table in doc.tables:
            rows: list[str] = []
            for i, row in enumerate(table.rows):
                cells = " | ".join(c.text.strip() for c in row.cells)
                rows.append(f"| {cells} |")
                if i == 0:
                    rows.append("|" + " --- |" * len(row.cells))
            parts.append("\n".join(rows))
        return "\n\n".join(parts)

    # ── Excel (.xlsx) ─────────────────────────────────────────────────────────

    def _create_xlsx(self, path: Path, title: str, content: str, styles: dict) -> None:
        import openpyxl  # type: ignore[import]
        from openpyxl.styles import Font  # type: ignore[import]
        wb = openpyxl.Workbook()
        ws = wb.active
        ws.title = title[:31] if title else "Sheet1"  # Excel max 31 chars
        # Parse CSV-like or Markdown table content
        for row_idx, line in enumerate(content.split("\n"), start=1):
            if "|" in line:
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                for col_idx, val in enumerate(cells, start=1):
                    cell = ws.cell(row=row_idx, column=col_idx, value=val)
                    if row_idx == 1 and styles.get("bold"):
                        cell.font = Font(bold=True)
                    fmt = styles.get("cell_format")
                    if fmt:
                        cell.number_format = fmt
            elif line.strip():
                ws.cell(row=row_idx, column=1, value=line.strip())
        wb.save(path)

    def _append_xlsx(self, path: Path, content: str) -> None:
        import openpyxl  # type: ignore[import]
        wb = openpyxl.load_workbook(path)
        ws = wb.active
        start_row = ws.max_row + 1
        for row_idx, line in enumerate(content.split("\n"), start=start_row):
            if "|" in line:
                cells = [c.strip() for c in line.strip().strip("|").split("|")]
                for col_idx, val in enumerate(cells, start=1):
                    ws.cell(row=row_idx, column=col_idx, value=val)
            elif line.strip():
                ws.cell(row=row_idx, column=1, value=line.strip())
        wb.save(path)

    def _read_xlsx(self, path: Path) -> str:
        import openpyxl  # type: ignore[import]
        wb = openpyxl.load_workbook(path, read_only=True, data_only=True)
        parts: list[str] = []
        for ws in wb.worksheets:
            parts.append(f"## Sheet: {ws.title}")
            rows_data: list[list[str]] = []
            for row in ws.iter_rows(values_only=True):
                rows_data.append([str(c) if c is not None else "" for c in row])
            if rows_data:
                # Header row separator
                header = "| " + " | ".join(rows_data[0]) + " |"
                sep = "|" + " --- |" * len(rows_data[0])
                parts.append(header)
                parts.append(sep)
                for row in rows_data[1:]:
                    parts.append("| " + " | ".join(row) + " |")
        return "\n\n".join(parts)

    # ── PowerPoint (.pptx) ────────────────────────────────────────────────────

    def _create_pptx(self, path: Path, title: str, content: str, styles: dict) -> None:
        from pptx import Presentation  # type: ignore[import]
        from pptx.util import Inches, Pt  # type: ignore[import]
        prs = Presentation()
        # Title slide
        slide_layout = prs.slide_layouts[0]
        slide = prs.slides.add_slide(slide_layout)
        slide.shapes.title.text = title
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text = content[:200] if content else ""

        # Additional slides from content sections (split by ---)
        sections = content.split("---")
        for section in sections[1:]:
            lines = [l for l in section.strip().split("\n") if l.strip()]
            if not lines:
                continue
            content_layout = prs.slide_layouts[1]
            s = prs.slides.add_slide(content_layout)
            s.shapes.title.text = lines[0].lstrip("#").strip()
            body = "\n".join(lines[1:])
            if s.placeholders and len(s.placeholders) > 1:
                tf = s.placeholders[1].text_frame
                tf.text = body
                if styles.get("font_size"):
                    for para in tf.paragraphs:
                        for run in para.runs:
                            run.font.size = Pt(int(styles["font_size"]))
        prs.save(path)

    def _append_pptx(self, path: Path, content: str) -> None:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(path)
        lines = [l for l in content.strip().split("\n") if l.strip()]
        if not lines:
            return
        layout = prs.slide_layouts[1]
        slide = prs.slides.add_slide(layout)
        slide.shapes.title.text = lines[0].lstrip("#").strip()
        body = "\n".join(lines[1:])
        if slide.placeholders and len(slide.placeholders) > 1:
            slide.placeholders[1].text_frame.text = body
        prs.save(path)

    def _read_pptx(self, path: Path) -> str:
        from pptx import Presentation  # type: ignore[import]
        prs = Presentation(path)
        parts: list[str] = []
        for i, slide in enumerate(prs.slides, start=1):
            parts.append(f"## Slide {i}")
            for shape in slide.shapes:
                if shape.has_text_frame:
                    for para in shape.text_frame.paragraphs:
                        text = para.text.strip()
                        if text:
                            parts.append(text)
        return "\n\n".join(parts)

    # ── Schema ────────────────────────────────────────────────────────────────

    def get_schema(self) -> dict[str, Any]:
        return {
            "name": self.name,
            "description": self.description,
            "parameters": {
                "type": "object",
                "properties": {
                    "action": {
                        "type": "string",
                        "enum": list(_VALID_ACTIONS),
                        "description": "操作类型",
                    },
                    "path": {
                        "type": "string",
                        "description": "文件路径（相对于 WORKSPACE_DIR），含扩展名，如 reports/monthly.xlsx",
                    },
                    "title": {"type": "string", "description": "文档标题（create 时使用）"},
                    "content": {
                        "type": "string",
                        "description": "文档内容（Markdown 格式，支持 # 标题、| 表格、--- 分隔）",
                    },
                    "styles": {
                        "type": "object",
                        "description": "样式选项",
                        "properties": {
                            "bold": {"type": "boolean"},
                            "italic": {"type": "boolean"},
                            "font_size": {"type": "integer", "description": "字号（磅）"},
                            "heading_level": {"type": "integer", "description": "标题级别 1-6"},
                            "cell_format": {"type": "string", "description": "Excel 单元格格式，如 '#,##0.00'"},
                        },
                    },
                },
                "required": ["action", "path"],
            },
        }


__all__ = ["OfficeTool"]
